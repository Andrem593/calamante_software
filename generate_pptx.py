import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Configurar tamaño widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Definir colores corporativos
    COLOR_BG_DARK = RGBColor(15, 23, 42)      # Slate 900 (#0F172A)
    COLOR_CARD_BG = RGBColor(30, 41, 59)      # Slate 800 (#1E293B)
    COLOR_PRIMARY = RGBColor(3, 169, 244)     # Sky Blue (#03A9F4)
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252) # Slate 50 (#F8FAFC)
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# Slate 400 (#94A3B8)

    blank_slide_layout = prs.slide_layouts[6] # Completamente en blanco para diseño libre

    # Carpeta del proyecto para capturas reales
    capturas_dir = "/Users/mac/calamante/calamante_software/assets/images/capturas_reales"
    
    # Mapeo de capturas reales (profile.png es el login real!)
    img_login = os.path.join(capturas_dir, "profile.png")
    img_dash = os.path.join(capturas_dir, "dashboard.png")
    img_history = os.path.join(capturas_dir, "history.png")
    img_detail_1 = os.path.join(capturas_dir, "history_datail.png") # spelling in dir
    img_detail_2 = os.path.join(capturas_dir, "history_detail_2.png")
    img_profile = os.path.join(capturas_dir, "my_profile.png")
    img_order_1 = os.path.join(capturas_dir, "order.png")
    img_order_2 = os.path.join(capturas_dir, "order_2.png")
    img_sig_1 = os.path.join(capturas_dir, "signature.png")
    img_sig_2 = os.path.join(capturas_dir, "signature_2.png")

    mockup_login = "/Users/mac/.gemini/antigravity-cli/brain/dcd76b58-ae79-4b81-8a73-fbef3e2ad6b7/login_screen_mockup_1779932633335.png"

    # Resolver fallback de login
    login_path = img_login if os.path.exists(img_login) else mockup_login

    print("\n--- EVALUANDO IMÁGENES PARA LA PRESENTACIÓN ---")
    print(f"Login (profile.png): {'[REAL]' if os.path.exists(img_login) else '[MOCKUP FALLBACK]'}")
    print(f"Dashboard:           {'[REAL]' if os.path.exists(img_dash) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print(f"Historial:           {'[REAL]' if os.path.exists(img_history) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print(f"Pedido 1:            {'[REAL]' if os.path.exists(img_order_1) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print(f"Pedido 2:            {'[REAL]' if os.path.exists(img_order_2) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print(f"Firma 1:             {'[REAL]' if os.path.exists(img_sig_1) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print(f"Firma 2:             {'[REAL]' if os.path.exists(img_sig_2) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print(f"Detalle 1:           {'[REAL]' if os.path.exists(img_detail_1) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print(f"Detalle 2:           {'[REAL]' if os.path.exists(img_detail_2) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print(f"Perfil:              {'[REAL]' if os.path.exists(img_profile) else '[MOCKUP FALLBACK (NO DISPONIBLE)]'}")
    print("-----------------------------------------------\n")

    # Helper para aplicar fondo oscuro a una diapositiva
    def apply_dark_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG_DARK

    # Helper para añadir encabezado estándar a diapositiva
    def add_standard_header(slide, title_text):
        # Título
        tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Outfit"
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN

        # Marca esquina superior derecha
        brand_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.4), Inches(3.2), Inches(0.8))
        btf = brand_box.text_frame
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.RIGHT
        bp.text = "DADOS APP • Capacitación"
        bp.font.name = "Outfit"
        bp.font.size = Pt(12)
        bp.font.color.rgb = COLOR_PRIMARY

    # Helper para diapositivas con UN mockup
    def build_single_mockup_slide(slide, title, desc, img_path, bullets):
        apply_dark_bg(slide)
        add_standard_header(slide, title)

        # 1. Añadir mockup en el lado izquierdo
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.3), Inches(3.1), Inches(5.6))
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.3), Inches(3.1), Inches(5.6))
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CARD_BG
            shape.line.color.rgb = COLOR_PRIMARY
            tf = shape.text_frame
            tf.text = "[Captura no encontrada]"

        # 2. Descripción general
        desc_box = slide.shapes.add_textbox(Inches(4.1), Inches(1.3), Inches(8.6), Inches(0.6))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.name = "Inter"
        dp.font.size = Pt(16)
        dp.font.color.rgb = COLOR_TEXT_MUTED

        # 3. Bullets estructurados
        top_offset = 2.0
        for bullet in bullets:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.1), Inches(top_offset), Inches(8.6), Inches(1.1))
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_CARD_BG
            card.line.color.rgb = RGBColor(51, 65, 85) # Slate 700
            card.line.width = Pt(1)

            ctf = card.text_frame
            ctf.word_wrap = True
            ctf.margin_left = ctf.margin_right = Inches(0.2)
            ctf.margin_top = Inches(0.12)

            p1 = ctf.paragraphs[0]
            p1.text = f"✓  {bullet['title']}"
            p1.font.name = "Outfit"
            p1.font.size = Pt(15)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_PRIMARY

            p2 = ctf.add_paragraph()
            p2.text = bullet['desc']
            p2.font.name = "Inter"
            p2.font.size = Pt(12)
            p2.font.color.rgb = COLOR_TEXT_MUTED
            
            top_offset += 1.3

    # Helper para diapositivas con DOS mockups lado a lado
    def build_double_mockup_slide(slide, title, desc, img_path1, img_path2, label1, label2, bullets):
        apply_dark_bg(slide)
        add_standard_header(slide, title)

        # 1. Mockup 1 (Izquierdo)
        if os.path.exists(img_path1):
            slide.shapes.add_picture(img_path1, Inches(0.6), Inches(1.3), Inches(2.7), Inches(5.5))
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(2.7), Inches(5.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CARD_BG
            tf = shape.text_frame
            tf.text = f"[Captura {label1} no encontrada]"

        # Etiqueta Mockup 1
        lbl_box1 = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(2.7), Inches(0.3))
        ltf1 = lbl_box1.text_frame
        lp1 = ltf1.paragraphs[0]
        lp1.alignment = PP_ALIGN.CENTER
        lp1.text = label1
        lp1.font.name = "Inter"
        lp1.font.size = Pt(10)
        lp1.font.bold = True
        lp1.font.color.rgb = COLOR_PRIMARY

        # 2. Mockup 2 (Central)
        if os.path.exists(img_path2):
            slide.shapes.add_picture(img_path2, Inches(3.5), Inches(1.3), Inches(2.7), Inches(5.5))
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(1.3), Inches(2.7), Inches(5.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CARD_BG
            tf = shape.text_frame
            tf.text = f"[Captura {label2} no encontrada]"

        # Etiqueta Mockup 2
        lbl_box2 = slide.shapes.add_textbox(Inches(3.5), Inches(6.8), Inches(2.7), Inches(0.3))
        ltf2 = lbl_box2.text_frame
        lp2 = ltf2.paragraphs[0]
        lp2.alignment = PP_ALIGN.CENTER
        lp2.text = label2
        lp2.font.name = "Inter"
        lp2.font.size = Pt(10)
        lp2.font.bold = True
        lp2.font.color.rgb = COLOR_PRIMARY

        # 3. Descripción general (Lado derecho)
        desc_box = slide.shapes.add_textbox(Inches(6.4), Inches(1.3), Inches(6.3), Inches(0.6))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.name = "Inter"
        dp.font.size = Pt(15)
        dp.font.color.rgb = COLOR_TEXT_MUTED

        # 4. Bullets estructurados en tarjetas (Lado derecho)
        top_offset = 2.0
        for bullet in bullets:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(top_offset), Inches(6.3), Inches(1.1))
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_CARD_BG
            card.line.color.rgb = RGBColor(51, 65, 85)
            card.line.width = Pt(1)

            ctf = card.text_frame
            ctf.word_wrap = True
            ctf.margin_left = ctf.margin_right = Inches(0.2)
            ctf.margin_top = Inches(0.12)

            p1 = ctf.paragraphs[0]
            p1.text = f"✓  {bullet['title']}"
            p1.font.name = "Outfit"
            p1.font.size = Pt(14)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_PRIMARY

            p2 = ctf.add_paragraph()
            p2.text = bullet['desc']
            p2.font.name = "Inter"
            p2.font.size = Pt(11)
            p2.font.color.rgb = COLOR_TEXT_MUTED
            
            top_offset += 1.3

    # ----------------- SLIDE 1: PORTADA -----------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide1)

    logo_box = slide1.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.0), Inches(1.5))
    ltf = logo_box.text_frame
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lp.text = "❄️"
    lp.font.name = "Outfit"
    lp.font.size = Pt(64)
    lp.font.color.rgb = COLOR_PRIMARY

    title_box = slide1.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12.0), Inches(1.2))
    ttf = title_box.text_frame
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tp.text = "Dados App"
    tp.font.name = "Outfit"
    tp.font.size = Pt(60)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_TEXT_MAIN

    sub_box = slide1.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(12.0), Inches(0.8))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.text = "Toma de Pedidos y Facturación Móvil de Hielo Purificado"
    sp.font.name = "Outfit"
    sp.font.size = Pt(22)
    sp.font.color.rgb = COLOR_PRIMARY

    foot_box = slide1.shapes.add_textbox(Inches(0.6), Inches(5.3), Inches(12.0), Inches(0.8))
    ftf = foot_box.text_frame
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fp.text = "Presentación Oficial de Capacitación para Vendedores y Repartidores"
    fp.font.name = "Inter"
    fp.font.size = Pt(14)
    fp.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------- SLIDE 2: ACCESO Y SEGURIDAD (LOGIN) -----------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    build_single_mockup_slide(
        slide2,
        "Acceso Seguro e Instantáneo 🔐",
        "Garantiza el ingreso ágil del personal de ventas y logística en campo sin contraseñas engorrosas.",
        login_path,
        [
            {"title": "Autenticación Biométrica Integrada", "desc": "Permite iniciar sesión rápidamente usando la huella digital o el reconocimiento facial (FaceID/TouchID) del dispositivo."},
            {"title": "Auto-Login Inteligente", "desc": "La app recuerda la sesión y realiza un auto-login rápido de 5 segundos para que los usuarios no pierdan tiempo al abrir la app."},
            {"title": "Limpieza de Seguridad", "desc": "Al cerrar sesión, la base local se limpia de borradores y tokens para una protección integral frente a extravíos."}
        ]
    )

    # ----------------- SLIDE 3: DASHBOARD -----------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    build_single_mockup_slide(
        slide3,
        "Panel de Control Centralizado (Dashboard) 📊",
        "Toda la información del día disponible de un vistazo para el repartidor y el vendedor de ruta.",
        img_dash,
        [
            {"title": "Métricas de Ruta e Indicadores de Ventas", "desc": "Visualiza al instante pedidos 'En proceso', 'Pedidos del Mes' y el total acumulado en 'Ventas del Mes'."},
            {"title": "Módulo Especial: Entregas de Hoy 🚚", "desc": "Tarjeta en degradé que destaca el total de entregas pendientes para hoy. El botón 'Ver Detalles' permite planificar la ruta de inmediato."},
            {"title": "Buscador Rápido e Historial", "desc": "Acceso inmediato al catálogo de stock de productos de hielo y al listado con los pedidos realizados recientemente."}
        ]
    )

    # ----------------- SLIDE 4: TOMA DE PEDIDO (DOBLE) -----------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    build_double_mockup_slide(
        slide4,
        "Proceso de Toma de Pedido 📝",
        "El formulario ágil y resiliente que los vendedores usan en campo para levantar órdenes.",
        img_order_1,
        img_order_2,
        "Parte 1: Catálogo",
        "Parte 2: GPS y Datos",
        [
            {"title": "Borradores Automáticos (Cero pérdidas)", "desc": "La app auto-guarda todos los datos ingresados en tiempo real. Si el teléfono se apaga o la señal falla, al abrir la app todo se recupera intacto."},
            {"title": "Ubicación Georreferenciada por GPS", "desc": "Usa el GPS para guardar la coordenada al crear el pedido, visualizándolo en un mapa interactivo para auditoría de rutas."},
            {"title": "Filtros por Categoría y Resumen Financiero", "desc": "Chips para filtrar productos, con cálculo instantáneo de Subtotal, IVA (15%) y Total en la barra inferior."}
        ]
    )

    # ----------------- SLIDE 5: REVISION Y FIRMA DIGITAL (DOBLE) -----------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    build_double_mockup_slide(
        slide5,
        "Revisión y Aceptación con Firma Digital ✍️",
        "Cierre legal y de conformidad en cada entrega, digitalizando la firma física directamente.",
        img_sig_1,
        img_sig_2,
        "Parte 1: Totales",
        "Parte 2: Firma Táctil",
        [
            {"title": "Desglose Contable Detallado", "desc": "Muestra al cliente el detalle exacto de ítems, el subtotal, el cálculo preciso de IVA y el Total General antes de la firma."},
            {"title": "Identificación de Responsable Físico", "desc": "Campos obligatorios para capturar el Nombre de quien recibe y su RUC/Cédula, previniendo disputas contables o entregas perdidas."},
            {"title": "Firma Táctil Digitalizada", "desc": "Lienzo digital interactivo para firmar con el dedo sobre la pantalla móvil. Se limpia fácilmente para reintentos y viaja inmutable a la base de datos Laravel."}
        ]
    )

    # ----------------- SLIDE 6: DETALLE DE PEDIDO Y BOTON DE ENTREGA (DOBLE) -----------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    build_double_mockup_slide(
        slide6,
        "Detalle del Pedido y Despacho 🚚",
        "Consola de ruta del repartidor para auditar las entregas y resguardar la facturación.",
        img_detail_1,
        img_detail_2,
        "Parte 1: Timeline",
        "Parte 2: Despacho",
        [
            {"title": "Timeline de Progreso Logístico", "desc": "Muestra visualmente la etapa de la orden: Recibido ➔ Facturado ➔ En camino (automático para hoy) ➔ Entregado."},
            {"title": "Auditoría de Firma e Identificación", "desc": "Recupera la imagen de la firma de conformidad y los datos del responsable que fueron capturados en campo."},
            {"title": "Alerta Inteligente de Facturación 🚨", "desc": "Si el repartidor presiona 'Entregar Pedido' sin el switch 'Facturado' activo, la app bloquea temporalmente el flujo y advierte al chofer para evitar entregar hielo sin factura."}
        ]
    )

    # ----------------- SLIDE 7: HISTORIAL DE PEDIDOS -----------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    build_single_mockup_slide(
        slide7,
        "Historial de Pedidos e Histórico 🔍",
        "Módulo centralizado para auditar e inspeccionar todo el histórico de transacciones.",
        img_history,
        [
            {"title": "Filtros Rápidos por Categoría de Estado", "desc": "Chips de selección instantánea para agrupar las órdenes por: Todos, Entregas Hoy, Pendientes, Completados o Cancelados."},
            {"title": "Buscador Global e Integrado", "desc": "Permite ubicar cualquier orden digitando directamente el ID (#) de pedido o el nombre comercial del cliente."},
            {"title": "Badges de Facturación de un Vistazo", "desc": "Cada tarjeta detalla visiblemente si la orden está en estado 'Facturado' o 'Prefacturado', simplificando el control contable en ruta."}
        ]
    )

    # ----------------- SLIDE 8: PERFIL Y SOPORTE -----------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    build_single_mockup_slide(
        slide8,
        "Perfil de Usuario y Soporte Técnico 👤",
        "Área de control individual, gestión de seguridad y soporte corporativo.",
        img_profile,
        [
            {"title": "Estadísticas Mensuales de Desempeño", "desc": "Consolida las ventas individuales acumuladas en dólares y los pedidos completados en el mes en curso."},
            {"title": "Canal de Soporte Directo (WhatsApp)", "desc": "Redirección rápida a WhatsApp del soporte corporativo con plantilla de mensaje lista para asistencia técnica en ruta."},
            {"title": "Administración Segura de Contraseña", "desc": "Diálogo integrado que permite el cambio de credenciales de seguridad directo con el servidor sin salir de la app."}
        ]
    )

    # ----------------- SLIDE 9: CAPACITACION Y RESUMEN -----------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide9)
    add_standard_header(slide9, "Resumen de Capacitación Comercial")

    res_desc = slide9.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.5))
    rtf = res_desc.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = "Conceptos esenciales que el vendedor y el repartidor deben dominar al 100% mañana:"
    rp.font.name = "Inter"
    rp.font.size = Pt(15)
    rp.font.color.rgb = COLOR_TEXT_MUTED

    summary_data = [
        {"icon": "⚡", "title": "Borrador Auto-Save", "desc": "El autoguardado en base local previene perder información de pedidos si el celular se descarga o falla la señal de ruta en campo."},
        {"icon": "📍", "title": "GPS e Integración", "desc": "La toma de ubicación en el mapa al crear el pedido garantiza a la administración que el vendedor visitó físicamente al cliente."},
        {"icon": "✍️", "title": "Firma Digital Táctil", "desc": "El cliente firma con el dedo directamente sobre la pantalla móvil, dando conformidad legal e instantánea a la entrega del hielo."},
        {"icon": "🛡️", "title": "Alerta de Facturación", "desc": "El botón de entrega bloquea temporalmente el despacho si el hielo no ha sido facturado oficialmente, reduciendo fugas de cobros."}
    ]

    left_pos = 0.6
    for item in summary_data:
        box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.9), Inches(2.8), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD_BG
        box.line.color.rgb = RGBColor(55, 65, 81)
        
        btf = box.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_right = Inches(0.2)
        btf.margin_top = Inches(0.3)
        
        p_icon = btf.paragraphs[0]
        p_icon.alignment = PP_ALIGN.CENTER
        p_icon.text = item["icon"]
        p_icon.font.size = Pt(36)
        
        p_title = btf.add_paragraph()
        p_title.alignment = PP_ALIGN.CENTER
        p_title.text = f"\n{item['title']}"
        p_title.font.name = "Outfit"
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        
        p_desc = btf.add_paragraph()
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.text = f"\n{item['desc']}"
        p_desc.font.name = "Inter"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        
        left_pos += 3.05

    output_path = "/Users/mac/calamante/calamante_software/presentacion_dados_app.pptx"
    prs.save(output_path)
    print(f"Presentación PPTX actualizada exitosamente en: {output_path}")

if __name__ == "__main__":
    create_presentation()
